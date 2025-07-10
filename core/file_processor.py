import fitz  # PyMuPDF
import docx
import csv
import pptx
import os

def extract_text(file_path, filetype):
    if filetype == "pdf":
        text = ""
        doc = fitz.open(file_path)
        for page in doc:
            text += page.get_text()
        return text

    elif filetype == "docx":
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])

    elif filetype == "pptx":
        presentation = pptx.Presentation(file_path)
        text_runs = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)

    elif filetype == "csv":
        with open(file_path, newline='', encoding='utf-8') as csvfile:
            return csvfile.read()

    elif filetype in ["txt", "md"]:
        with open(file_path, encoding='utf-8') as f:
            return f.read()

    else:
        raise ValueError(f"Unsupported file type: {filetype}")
def chunk_text(text, max_tokens=500, overlap=50):
    """
    Splits the input text into overlapping chunks based on max_tokens.
    Assumes ~1 token per 4 characters as a rough estimate.
    """
    approx_chars_per_token = 4
    max_chars = max_tokens * approx_chars_per_token
    overlap_chars = overlap * approx_chars_per_token

    chunks = []
    start = 0
    while start < len(text):
        end = min(start + max_chars, len(text))
        chunks.append(text[start:end])
        start += max_chars - overlap_chars

    return chunks

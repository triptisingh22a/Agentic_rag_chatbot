import pdfplumber
import os
import docx
import pandas as pd
from pptx import Presentation

def parse_pdf(file_path):
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() + "\n"
    return [t.strip() for t in text.split("\n") if t.strip()]

def parse_docx(file_path):
    doc = docx.Document(file_path)
    return [para.text for para in doc.paragraphs if para.text.strip()]

def parse_csv(file_path):
    df = pd.read_csv(file_path)
    return [row for row in df.astype(str).agg(' | '.join, axis=1)]

def parse_pptx(file_path):
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return [t.strip() for t in text_runs if t.strip()]

def parse_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return [line.strip() for line in f.readlines() if line.strip()]

def parse_document(file):
    filename = file.name
    file_path = os.path.join("data", "uploaded_files", filename)
    with open(file_path, "wb") as f:
        f.write(file.read())

    if filename.endswith(".pdf"):
        return parse_pdf(file_path)
    elif filename.endswith(".docx"):
        return parse_docx(file_path)
    elif filename.endswith(".csv"):
        return parse_csv(file_path)
    elif filename.endswith(".pptx"):
        return parse_pptx(file_path)
    elif filename.endswith(".txt") or filename.endswith(".md"):
        return parse_txt(file_path)
    else:
        return []
